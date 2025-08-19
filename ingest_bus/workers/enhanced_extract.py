"""
Enhanced Extract Worker Module for TORI Ingest Bus
Adds support for all file types: PDF, DOCX, CSV, PPTX, XLSX, JSON, MD, TXT

Integrates with:
- ConceptMesh (concept graph)
- Ghost Collective (AI personas) 
- ψMesh (semantic associations)
- ScholarSphere (archival)
"""

import os
import sys
import json
import hashlib
import logging
import asyncio
import pandas as pd
from io import BytesIO
from typing import Dict, List, Any, Optional, Union, Tuple
from datetime import datetime
from pathlib import Path

# Document parsing libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

# Import existing models
from models.schemas import (
    IngestStatus, DocumentType, FailureCode,
    Chunk, ConceptVectorLink, IngestJob
)

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("tori-ingest.enhanced_extract")

# Load configuration
try:
    config_path = Path(__file__).parent.parent.parent / "conversation_config.json"
    with open(config_path, "r") as f:
        config = json.load(f)
except Exception as e:
    logger.warning(f"Could not load configuration: {str(e)}")
    config = {
        "scholar_sphere": {
            "enabled": True,
            "encoder_version": "v2.5.0",
            "chunk_size": 512,
            "chunk_overlap": 128,
            "max_concepts_per_chunk": 12
        },
        "concept_mesh": {
            "enabled": True,
            "api_endpoint": "http://localhost:8081"
        },
        "psi_mesh": {
            "enabled": True,
            "verification_threshold": 0.85
        }
    }

class DocumentProcessor:
    """Enhanced document processor supporting all file types"""
    
    def __init__(self):
        self.supported_types = {
            'pdf': self.extract_pdf,
            'docx': self.extract_docx,
            'doc': self.extract_docx,
            'csv': self.extract_csv,
            'pptx': self.extract_pptx,
            'xlsx': self.extract_xlsx,
            'json': self.extract_json,
            'txt': self.extract_text,
            'md': self.extract_markdown
        }
    
    async def process_document(self, file_path: Optional[str], file_content: Optional[bytes], 
                             file_type: str, job: IngestJob) -> Optional[Dict[str, Any]]:
        """
        Process document and extract structured content
        
        Returns:
            Dict containing:
            - text: extracted text content
            - metadata: document metadata  
            - structure: document structure (headings, sections, etc.)
            - concepts: extracted semantic concepts
        """
        logger.info(f"Processing {file_type} document: job_id={job.id}")
        
        if file_type.lower() not in self.supported_types:
            logger.error(f"Unsupported file type: {file_type}")
            return None
        
        try:
            processor = self.supported_types[file_type.lower()]
            result = await processor(file_path, file_content, job)
            
            if result:
                # Add common metadata
                result['metadata']['file_type'] = file_type
                result['metadata']['processed_at'] = datetime.now().isoformat()
                result['metadata']['job_id'] = job.id
                
                # Extract semantic concepts
                if result.get('text'):
                    concepts = await self.extract_semantic_concepts(result['text'], job)
                    result['concepts'] = concepts
                
                logger.info(f"Successfully processed {file_type} document")
                return result
            
        except Exception as e:
            logger.exception(f"Error processing {file_type} document: {str(e)}")
            return None
    
    async def extract_pdf(self, file_path: Optional[str], file_content: Optional[bytes], 
                         job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from PDF"""
        if not PdfReader:
            logger.error("PyPDF2 not available for PDF extraction")
            return None
        
        try:
            if file_path:
                reader = PdfReader(file_path)
            else:
                reader = PdfReader(BytesIO(file_content))
            
            text = ""
            structure = []
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                text += page_text
                
                # Try to detect headings by font size (basic approach)
                structure.append({
                    'type': 'page',
                    'number': page_num + 1,
                    'text_length': len(page_text)
                })
            
            metadata = {
                'pages': len(reader.pages),
                'pdf_metadata': reader.metadata if hasattr(reader, 'metadata') else {}
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure
            }
            
        except Exception as e:
            logger.exception(f"Error extracting PDF: {str(e)}")
            return None
    
    async def extract_docx(self, file_path: Optional[str], file_content: Optional[bytes], 
                          job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from DOCX/DOC files"""
        if not Document:
            logger.error("python-docx not available for DOCX extraction")
            return None
        
        try:
            if file_path:
                doc = Document(file_path)
            else:
                doc = Document(BytesIO(file_content))
            
            text = ""
            structure = []
            
            for para in doc.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    text += para_text + "\n\n"
                    
                    # Detect headings by style
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                        structure.append({
                            'type': 'heading',
                            'level': level,
                            'text': para_text
                        })
            
            # Extract tables
            tables_text = ""
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    tables_text += row_text + "\n"
            
            if tables_text:
                text += "\n\nTables:\n" + tables_text
            
            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'core_properties': {
                    'title': doc.core_properties.title,
                    'author': doc.core_properties.author,
                    'subject': doc.core_properties.subject
                } if hasattr(doc, 'core_properties') else {}
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure
            }
            
        except Exception as e:
            logger.exception(f"Error extracting DOCX: {str(e)}")
            return None
    
    async def extract_csv(self, file_path: Optional[str], file_content: Optional[bytes], 
                         job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from CSV files"""
        try:
            if file_path:
                df = pd.read_csv(file_path)
            else:
                df = pd.read_csv(BytesIO(file_content))
            
            # Convert to text representation
            text = f"Dataset with {len(df)} rows and {len(df.columns)} columns.\n\n"
            text += f"Columns: {', '.join(df.columns)}\n\n"
            
            # Add sample data
            text += "Sample data:\n"
            text += df.head(10).to_string(index=False)
            
            # Add statistical summary for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                text += "\n\nStatistical Summary:\n"
                text += df[numeric_cols].describe().to_string()
            
            structure = [
                {
                    'type': 'dataset',
                    'rows': len(df),
                    'columns': list(df.columns),
                    'data_types': {col: str(dtype) for col, dtype in df.dtypes.items()}
                }
            ]
            
            metadata = {
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': list(df.columns),
                'data_types': {col: str(dtype) for col, dtype in df.dtypes.items()},
                'memory_usage': df.memory_usage(deep=True).sum()
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure,
                'raw_data': df.to_dict('records')[:100]  # First 100 rows
            }
            
        except Exception as e:
            logger.exception(f"Error extracting CSV: {str(e)}")
            return None
    
    async def extract_pptx(self, file_path: Optional[str], file_content: Optional[bytes], 
                          job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from PPTX files"""
        if not Presentation:
            logger.error("python-pptx not available for PPTX extraction")
            return None
        
        try:
            if file_path:
                prs = Presentation(file_path)
            else:
                prs = Presentation(BytesIO(file_content))
            
            text = ""
            structure = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                text += slide_text + "\n\n"
                
                structure.append({
                    'type': 'slide',
                    'number': slide_num + 1,
                    'shapes': len(slide.shapes)
                })
            
            metadata = {
                'slides': len(prs.slides),
                'slide_layouts': len(prs.slide_layouts)
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure
            }
            
        except Exception as e:
            logger.exception(f"Error extracting PPTX: {str(e)}")
            return None
    
    async def extract_xlsx(self, file_path: Optional[str], file_content: Optional[bytes], 
                          job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from XLSX files"""
        if not openpyxl:
            logger.error("openpyxl not available for XLSX extraction")
            return None
        
        try:
            if file_path:
                workbook = openpyxl.load_workbook(file_path)
            else:
                workbook = openpyxl.load_workbook(BytesIO(file_content))
            
            text = f"Excel workbook with {len(workbook.sheetnames)} sheets.\n\n"
            structure = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Convert sheet to DataFrame for easier handling
                data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append(row)
                
                if data:
                    df = pd.DataFrame(data[1:], columns=data[0] if data else None)
                    
                    text += f"Sheet: {sheet_name}\n"
                    text += f"Dimensions: {sheet.max_row} rows x {sheet.max_column} columns\n"
                    
                    # Add sample data
                    if not df.empty:
                        text += "Sample data:\n"
                        text += df.head(5).to_string(index=False) + "\n\n"
                    
                    structure.append({
                        'type': 'sheet',
                        'name': sheet_name,
                        'rows': sheet.max_row,
                        'columns': sheet.max_column
                    })
            
            metadata = {
                'sheets': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure
            }
            
        except Exception as e:
            logger.exception(f"Error extracting XLSX: {str(e)}")
            return None
    
    async def extract_json(self, file_path: Optional[str], file_content: Optional[bytes], 
                          job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from JSON files"""
        try:
            if file_path:
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = json.loads(file_content.decode('utf-8'))
            
            # Convert JSON to readable text
            text = "JSON Document:\n\n"
            text += json.dumps(data, indent=2, ensure_ascii=False)
            
            # Extract structure information
            structure = self._analyze_json_structure(data)
            
            metadata = {
                'json_type': type(data).__name__,
                'size_estimate': len(str(data))
            }
            
            if isinstance(data, list):
                metadata['array_length'] = len(data)
            elif isinstance(data, dict):
                metadata['object_keys'] = list(data.keys())
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure,
                'raw_data': data
            }
            
        except Exception as e:
            logger.exception(f"Error extracting JSON: {str(e)}")
            return None
    
    async def extract_text(self, file_path: Optional[str], file_content: Optional[bytes], 
                          job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from plain text files"""
        try:
            if file_path:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            else:
                text = file_content.decode('utf-8')
            
            # Basic structure analysis
            lines = text.split('\n')
            structure = [{
                'type': 'text_document',
                'lines': len(lines),
                'characters': len(text)
            }]
            
            metadata = {
                'lines': len(lines),
                'characters': len(text),
                'words': len(text.split())
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure
            }
            
        except Exception as e:
            logger.exception(f"Error extracting text: {str(e)}")
            return None
    
    async def extract_markdown(self, file_path: Optional[str], file_content: Optional[bytes], 
                              job: IngestJob) -> Optional[Dict[str, Any]]:
        """Extract content from Markdown files"""
        try:
            if file_path:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            else:
                text = file_content.decode('utf-8')
            
            # Extract structure from markdown
            lines = text.split('\n')
            structure = []
            
            for line_num, line in enumerate(lines):
                line = line.strip()
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    heading_text = line.lstrip('# ').strip()
                    structure.append({
                        'type': 'heading',
                        'level': level,
                        'text': heading_text,
                        'line': line_num + 1
                    })
            
            metadata = {
                'lines': len(lines),
                'characters': len(text),
                'headings': len(structure)
            }
            
            return {
                'text': text,
                'metadata': metadata,
                'structure': structure
            }
            
        except Exception as e:
            logger.exception(f"Error extracting markdown: {str(e)}")
            return None
    
    def _analyze_json_structure(self, data, path="", max_depth=3) -> List[Dict]:
        """Recursively analyze JSON structure"""
        structure = []
        
        if max_depth <= 0:
            return structure
        
        if isinstance(data, dict):
            for key, value in data.items():
                current_path = f"{path}.{key}" if path else key
                structure.append({
                    'type': 'object_key',
                    'path': current_path,
                    'value_type': type(value).__name__
                })
                
                # Recurse into nested structures
                if isinstance(value, (dict, list)):
                    structure.extend(self._analyze_json_structure(value, current_path, max_depth - 1))
        
        elif isinstance(data, list) and data:
            # Analyze first few items
            for i, item in enumerate(data[:3]):
                current_path = f"{path}[{i}]" if path else f"[{i}]"
                structure.append({
                    'type': 'array_item',
                    'path': current_path,
                    'value_type': type(item).__name__
                })
                
                if isinstance(item, (dict, list)):
                    structure.extend(self._analyze_json_structure(item, current_path, max_depth - 1))
        
        return structure
    
    async def extract_semantic_concepts(self, text: str, job: IngestJob) -> List[Dict[str, Any]]:
        """
        Extract semantic concepts from text using TORI's advanced concept extraction
        """
        try:
            # Import existing concept extraction from ingest_pdf
            sys.path.append(str(Path(__file__).parent.parent.parent / "ingest_pdf"))
            
            from extract_blocks import extract_concept_blocks
            from features import build_feature_matrix
            from spectral import spectral_embed
            from clustering import run_oscillator_clustering
            from scoring import score_clusters
            from keywords import extract_keywords
            
            # Create temporary blocks from text
            # Split into paragraphs as blocks
            blocks = [block.strip() for block in text.split('\n\n') if block.strip()]
            
            if not blocks:
                return []
            
            # Use TORI's advanced concept extraction
            feats, vocab = build_feature_matrix(blocks)
            emb = spectral_embed(feats, k=16)
            labels = run_oscillator_clustering(emb)
            top_cids = score_clusters(labels, emb)[:8]  # Top 8 concepts
            
            concepts = []
            for cid in top_cids:
                mem = [i for i, l in enumerate(labels) if l == cid]
                cluster_blocks = [blocks[i] for i in mem]
                other_blocks = [blocks[i] for i in range(len(blocks)) if i not in mem]
                
                keywords = extract_keywords(cluster_blocks, other_blocks, n=3)
                concept_name = " ".join(w.capitalize() for w in keywords) or f"Concept-{cid+1}"
                
                concepts.append({
                    'name': concept_name,
                    'keywords': keywords,
                    'cluster_id': int(cid),
                    'supporting_blocks': mem,
                    'context': cluster_blocks[0] if cluster_blocks else ""
                })
            
            logger.info(f"Extracted {len(concepts)} semantic concepts")
            return concepts
            
        except Exception as e:
            logger.exception(f"Error extracting semantic concepts: {str(e)}")
            return []

# Global processor instance
processor = DocumentProcessor()

# Enhanced extraction functions for ingest-bus compatibility
async def extract_document(file_path: Optional[str], file_content: Optional[bytes], 
                          document_type: str, job: IngestJob) -> Optional[Dict[str, Any]]:
    """
    Main extraction function that routes to appropriate processor
    """
    return await processor.process_document(file_path, file_content, document_type, job)

# Maintain backwards compatibility with existing extract.py functions
async def extract_pdf(file_path: Optional[str], file_content: Optional[bytes], job: IngestJob) -> Optional[str]:
    """Backwards compatible PDF extraction"""
    result = await processor.extract_pdf(file_path, file_content, job)
    return result['text'] if result else None

async def extract_conversation(file_path: Optional[str], file_content: Optional[bytes], job: IngestJob) -> Optional[str]:
    """Backwards compatible conversation extraction - delegates to original implementation"""
    # Import original function
    from extract import extract_conversation as original_extract_conversation
    return await original_extract_conversation(file_path, file_content, job)

async def extract_text(file_path: Optional[str], file_content: Optional[bytes], job: IngestJob) -> Optional[str]:
    """Backwards compatible text extraction"""
    result = await processor.extract_text(file_path, file_content, job)
    return result['text'] if result else None

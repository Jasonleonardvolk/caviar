"""
Complete File Type Handlers for TORI Ingest Bus
Production-ready parsers for ALL supported document types

Using best-in-class libraries:
- mammoth for DOCX (better than python-docx for HTML conversion)
- python-pptx for PowerPoint
- openpyxl for Excel  
- pandas for CSV analysis
- PyPDF2 for PDF processing
"""

import mammoth
import json
import logging
import asyncio
import pandas as pd
import numpy as np
import re
from io import BytesIO, StringIO
from typing import Dict, List, Any, Optional, Tuple, Union, Set
from datetime import datetime
from pathlib import Path
from collections import Counter, defaultdict

# Document processing libraries
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import mammoth
    MAMMOTH_AVAILABLE = True
except ImportError:
    MAMMOTH_AVAILABLE = False

# Import TORI models
from models.schemas import IngestJob

logger = logging.getLogger("tori-ingest.file_handlers")

class ParsedPayload:
    """
    Standardized output format for all document types
    Ensures consistent interface across TORI systems
    """
    
    def __init__(self):
        self.document_id: str = ""
        self.file_hash: str = ""
        self.source_metadata: Dict[str, Any] = {}
        self.extracted_text: str = ""
        self.structure_elements: List[Dict[str, Any]] = []
        self.semantic_concepts: List[Dict[str, Any]] = []
        self.raw_segments: List[Dict[str, Any]] = []
        self.integrity_metadata: Dict[str, Any] = {}
        self.processing_timestamps: Dict[str, str] = {}
        
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for serialization"""
        return {
            'document_id': self.document_id,
            'file_hash': self.file_hash,
            'source_metadata': self.source_metadata,
            'extracted_text': self.extracted_text,
            'structure_elements': self.structure_elements,
            'semantic_concepts': self.semantic_concepts,
            'raw_segments': self.raw_segments,
            'integrity_metadata': self.integrity_metadata,
            'processing_timestamps': self.processing_timestamps
        }

class ProductionFileHandlers:
    """
    Production-grade file handlers for all supported document types
    Optimized for reliability, performance, and consistency
    """
    
    def __init__(self):
        self.supported_handlers = {
            'pdf': self.handle_pdf,
            'docx': self.handle_docx,
            'doc': self.handle_docx,  # Same handler
            'pptx': self.handle_pptx,
            'xlsx': self.handle_xlsx,
            'xls': self.handle_xlsx,  # Same handler
            'csv': self.handle_csv,
            'json': self.handle_json,
            'txt': self.handle_text,
            'md': self.handle_markdown,
            'markdown': self.handle_markdown
        }
        
        # Dynamic Concept Extraction Configuration
        self.concept_config = {
            'base_confidence_threshold': 0.6,
            'high_quality_threshold': 0.8,
            'minimum_concept_length': 2,
            'maximum_concept_length': 50,
            'enable_entity_extraction': True,
            'enable_keyword_clustering': True,
            'enable_heading_concepts': True,
            'enable_frequency_analysis': True,
            'similarity_threshold': 0.85,
            'max_concepts_per_document': 100,  # Soft limit, can be exceeded for high-quality concepts
            'quality_over_quantity': True
        }
        
        logger.info("Production File Handlers initialized with Dynamic Concept Extraction")
        logger.info(f"Available libraries: PDF={PDF_AVAILABLE}, DOCX={MAMMOTH_AVAILABLE}, PPTX={PPTX_AVAILABLE}, XLSX={OPENPYXL_AVAILABLE}")
        logger.info(f"Concept extraction config: {self.concept_config}")
    
    async def process_document(self, file_content: bytes, file_type: str, 
                             filename: str, job: IngestJob) -> ParsedPayload:
        """
        Process document and return standardized ParsedPayload
        
        Args:
            file_content: Raw file bytes
            file_type: Document type (pdf, docx, etc.)
            filename: Original filename
            job: Ingest job metadata
            
        Returns:
            ParsedPayload with all extracted information
        """
        start_time = datetime.now()
        
        # Initialize payload
        payload = ParsedPayload()
        payload.document_id = job.id
        payload.file_hash = self._calculate_hash(file_content)
        payload.source_metadata = {
            'filename': filename,
            'file_type': file_type,
            'file_size': len(file_content),
            'job_id': job.id
        }
        payload.processing_timestamps['started'] = start_time.isoformat()
        
        # Route to appropriate handler
        if file_type.lower() in self.supported_handlers:
            handler = self.supported_handlers[file_type.lower()]
            try:
                await handler(file_content, payload, job)
                payload.processing_timestamps['completed'] = datetime.now().isoformat()
                
                # Extract semantic concepts from text with dynamic system
                if payload.extracted_text:
                    await self._extract_semantic_concepts_dynamic(payload)
                
                # Generate integrity metadata
                payload.integrity_metadata = self._generate_integrity_metadata(payload)
                
                logger.info(f"Successfully processed {file_type} document: {len(payload.extracted_text)} chars, {len(payload.semantic_concepts)} concepts extracted")
                
            except Exception as e:
                logger.exception(f"Error processing {file_type}: {e}")
                payload.processing_timestamps['failed'] = datetime.now().isoformat()
                payload.integrity_metadata['processing_error'] = str(e)
        else:
            logger.error(f"Unsupported file type: {file_type}")
            payload.integrity_metadata['error'] = f"Unsupported file type: {file_type}"
        
        return payload
    
    async def handle_pdf(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle PDF documents with advanced structure detection"""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 not available for PDF processing")
        
        try:
            reader = PdfReader(BytesIO(file_content))
            
            # Extract metadata
            if reader.metadata:
                payload.source_metadata.update({
                    'pdf_title': reader.metadata.get('/Title', ''),
                    'pdf_author': reader.metadata.get('/Author', ''),
                    'pdf_subject': reader.metadata.get('/Subject', ''),
                    'pdf_creator': reader.metadata.get('/Creator', ''),
                    'pdf_producer': reader.metadata.get('/Producer', ''),
                    'pdf_creation_date': str(reader.metadata.get('/CreationDate', '')),
                    'pdf_modification_date': str(reader.metadata.get('/ModDate', ''))
                })
            
            payload.source_metadata['total_pages'] = len(reader.pages)
            
            # Extract text and structure
            full_text = ""
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                full_text += page_text + "\n\n"
                
                # Create page structure element
                payload.structure_elements.append({
                    'type': 'page',
                    'page_number': page_num + 1,
                    'text_length': len(page_text),
                    'start_char': len(full_text) - len(page_text) - 2,
                    'end_char': len(full_text) - 2
                })
                
                # Create segments for each page
                if page_text.strip():
                    payload.raw_segments.append({
                        'segment_id': f"page_{page_num + 1}",
                        'type': 'page_content',
                        'text': page_text,
                        'metadata': {
                            'page_number': page_num + 1,
                            'char_start': len(full_text) - len(page_text) - 2,
                            'char_end': len(full_text) - 2
                        }
                    })
            
            payload.extracted_text = full_text.strip()
            
        except Exception as e:
            logger.exception(f"PDF processing error: {e}")
            raise
    
    async def handle_docx(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle DOCX/DOC documents with mammoth for better HTML conversion"""
        if not MAMMOTH_AVAILABLE:
            raise ImportError("mammoth not available for DOCX processing")
        
        try:
            # Use mammoth for better text extraction
            result = mammoth.extract_raw_text(BytesIO(file_content))
            payload.extracted_text = result.value
            
            # Also get HTML version for structure detection
            html_result = mammoth.convert_to_html(BytesIO(file_content))
            html_content = html_result.value
            
            # Parse HTML for structure elements
            structure_elements = self._parse_html_structure(html_content)
            payload.structure_elements.extend(structure_elements)
            
            # Create segments based on paragraphs
            paragraphs = payload.extracted_text.split('\n\n')
            char_position = 0
            
            for i, paragraph in enumerate(paragraphs):
                if paragraph.strip():
                    payload.raw_segments.append({
                        'segment_id': f"paragraph_{i + 1}",
                        'type': 'paragraph',
                        'text': paragraph.strip(),
                        'metadata': {
                            'paragraph_number': i + 1,
                            'char_start': char_position,
                            'char_end': char_position + len(paragraph)
                        }
                    })
                char_position += len(paragraph) + 2  # Account for \n\n
            
            # Extract document properties if available
            payload.source_metadata.update({
                'extraction_method': 'mammoth',
                'html_conversion_warnings': len(html_result.messages),
                'paragraph_count': len([p for p in paragraphs if p.strip()])
            })
            
        except Exception as e:
            logger.exception(f"DOCX processing error: {e}")
            raise
    
    async def handle_pptx(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available for PPTX processing")
        
        try:
            presentation = Presentation(BytesIO(file_content))
            
            payload.source_metadata.update({
                'slide_count': len(presentation.slides),
                'slide_layouts': len(presentation.slide_layouts)
            })
            
            full_text = ""
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = f"=== Slide {slide_num + 1} ===\n"
                slide_content = ""
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += shape.text + "\n"
                
                slide_text += slide_content
                full_text += slide_text + "\n\n"
                
                # Create slide structure element
                payload.structure_elements.append({
                    'type': 'slide',
                    'slide_number': slide_num + 1,
                    'shape_count': len(slide.shapes),
                    'text_length': len(slide_content),
                    'start_char': len(full_text) - len(slide_text) - 2,
                    'end_char': len(full_text) - 2
                })
                
                # Create segment for slide
                if slide_content.strip():
                    payload.raw_segments.append({
                        'segment_id': f"slide_{slide_num + 1}",
                        'type': 'slide_content',
                        'text': slide_content.strip(),
                        'metadata': {
                            'slide_number': slide_num + 1,
                            'shape_count': len(slide.shapes),
                            'char_start': len(full_text) - len(slide_text) - 2,
                            'char_end': len(full_text) - 2
                        }
                    })
            
            payload.extracted_text = full_text.strip()
            
        except Exception as e:
            logger.exception(f"PPTX processing error: {e}")
            raise
    
    async def handle_xlsx(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle Excel spreadsheets with comprehensive data analysis"""
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl not available for XLSX processing")
        
        try:
            workbook = openpyxl.load_workbook(BytesIO(file_content), read_only=True)
            
            payload.source_metadata.update({
                'sheet_count': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames
            })
            
            full_text = f"Excel Workbook: {len(workbook.sheetnames)} sheets\n\n"
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Convert sheet to pandas DataFrame for analysis
                data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append(row)
                
                if not data:
                    continue
                
                # Create DataFrame
                try:
                    df = pd.DataFrame(data[1:], columns=data[0] if data else None)
                    
                    # Generate text summary
                    sheet_text = f"=== Sheet: {sheet_name} ===\n"
                    sheet_text += f"Dimensions: {len(df)} rows × {len(df.columns)} columns\n"
                    sheet_text += f"Columns: {', '.join(str(col) for col in df.columns)}\n\n"
                    
                    # Add sample data
                    if len(df) > 0:
                        sheet_text += "Sample data:\n"
                        sheet_text += df.head(5).to_string(index=False) + "\n\n"
                    
                    # Add statistical summary for numeric columns
                    numeric_cols = df.select_dtypes(include=[np.number]).columns
                    if len(numeric_cols) > 0:
                        sheet_text += "Statistical Summary:\n"
                        sheet_text += df[numeric_cols].describe().to_string() + "\n\n"
                    
                    full_text += sheet_text
                    
                    # Create sheet structure element
                    payload.structure_elements.append({
                        'type': 'sheet',
                        'sheet_name': sheet_name,
                        'rows': len(df),
                        'columns': len(df.columns),
                        'numeric_columns': len(numeric_cols),
                        'data_types': {col: str(dtype) for col, dtype in df.dtypes.items()},
                        'start_char': len(full_text) - len(sheet_text),
                        'end_char': len(full_text)
                    })
                    
                    # Create segment for sheet
                    payload.raw_segments.append({
                        'segment_id': f"sheet_{sheet_name}",
                        'type': 'sheet_data',
                        'text': sheet_text,
                        'metadata': {
                            'sheet_name': sheet_name,
                            'rows': len(df),
                            'columns': len(df.columns),
                            'data_sample': df.head(3).to_dict('records') if len(df) > 0 else []
                        }
                    })
                    
                except Exception as sheet_error:
                    logger.warning(f"Error processing sheet {sheet_name}: {sheet_error}")
                    continue
            
            payload.extracted_text = full_text.strip()
            
        except Exception as e:
            logger.exception(f"XLSX processing error: {e}")
            raise
    
    async def handle_csv(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle CSV files with advanced data analysis"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin1', 'cp1252']
            df = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    csv_text = file_content.decode(encoding)
                    df = pd.read_csv(StringIO(csv_text))
                    used_encoding = encoding
                    break
                except (UnicodeDecodeError, pd.errors.EmptyDataError):
                    continue
            
            if df is None:
                raise ValueError("Could not decode CSV file with any supported encoding")
            
            # Generate comprehensive text analysis
            text_content = f"CSV Dataset Analysis\n\n"
            text_content += f"Dataset: {len(df)} rows × {len(df.columns)} columns\n"
            text_content += f"Encoding: {used_encoding}\n"
            text_content += f"Columns: {', '.join(df.columns)}\n\n"
            
            # Data types analysis
            text_content += "Data Types:\n"
            for col, dtype in df.dtypes.items():
                text_content += f"  {col}: {dtype}\n"
            text_content += "\n"
            
            # Sample data
            text_content += "Sample Data (first 5 rows):\n"
            text_content += df.head().to_string(index=False) + "\n\n"
            
            # Statistical summary
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text_content += "Statistical Summary:\n"
                text_content += df[numeric_cols].describe().to_string() + "\n\n"
            
            # Missing values analysis
            missing_data = df.isnull().sum()
            if missing_data.sum() > 0:
                text_content += "Missing Values:\n"
                for col, missing in missing_data.items():
                    if missing > 0:
                        text_content += f"  {col}: {missing} ({missing/len(df)*100:.1f}%)\n"
                text_content += "\n"
            
            # Unique values for categorical columns
            categorical_cols = df.select_dtypes(include=['object']).columns
            if len(categorical_cols) > 0:
                text_content += "Categorical Data Summary:\n"
                for col in categorical_cols[:5]:  # Limit to first 5 categorical columns
                    unique_count = df[col].nunique()
                    text_content += f"  {col}: {unique_count} unique values\n"
                    if unique_count <= 10:
                        text_content += f"    Values: {', '.join(str(v) for v in df[col].unique()[:10])}\n"
                text_content += "\n"
            
            payload.extracted_text = text_content
            
            # Store metadata
            payload.source_metadata.update({
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': list(df.columns),
                'data_types': {col: str(dtype) for col, dtype in df.dtypes.items()},
                'numeric_columns': len(numeric_cols),
                'categorical_columns': len(categorical_cols),
                'missing_values_total': int(missing_data.sum()),
                'encoding_used': used_encoding
            })
            
            # Create structure elements for columns
            for i, col in enumerate(df.columns):
                payload.structure_elements.append({
                    'type': 'column',
                    'column_name': col,
                    'column_index': i,
                    'data_type': str(df[col].dtype),
                    'unique_values': int(df[col].nunique()),
                    'missing_values': int(df[col].isnull().sum())
                })
            
            # Create segments for data subsets
            chunk_size = min(1000, len(df))  # Process in chunks of 1000 rows
            for i in range(0, len(df), chunk_size):
                chunk = df.iloc[i:i+chunk_size]
                chunk_text = f"Data rows {i+1} to {min(i+chunk_size, len(df))}:\n"
                chunk_text += chunk.to_string(index=False)
                
                payload.raw_segments.append({
                    'segment_id': f"rows_{i+1}_to_{min(i+chunk_size, len(df))}",
                    'type': 'data_chunk',
                    'text': chunk_text,
                    'metadata': {
                        'start_row': i,
                        'end_row': min(i+chunk_size, len(df)),
                        'row_count': len(chunk)
                    }
                })
            
        except Exception as e:
            logger.exception(f"CSV processing error: {e}")
            raise
    
    async def handle_json(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle JSON files with comprehensive structure analysis"""
        try:
            # Decode JSON
            json_text = file_content.decode('utf-8')
            data = json.loads(json_text)
            
            # Generate human-readable text representation
            text_content = "JSON Document Analysis\n\n"
            text_content += f"Structure: {self._get_json_type_description(data)}\n\n"
            
            # Add formatted JSON content
            text_content += "Content:\n"
            text_content += json.dumps(data, indent=2, ensure_ascii=False)[:5000]  # Limit to 5000 chars
            if len(json_text) > 5000:
                text_content += "\n... (content truncated)"
            
            payload.extracted_text = text_content
            
            # Analyze JSON structure
            structure_analysis = self._analyze_json_structure(data)
            payload.source_metadata.update({
                'json_type': type(data).__name__,
                'structure_depth': structure_analysis['max_depth'],
                'total_keys': structure_analysis['total_keys'],
                'total_values': structure_analysis['total_values'],
                'data_types': structure_analysis['data_types'],
                'file_size_chars': len(json_text)
            })
            
            # Create structure elements from JSON analysis
            payload.structure_elements.extend(structure_analysis['elements'])
            
            # Create segments for different parts of JSON
            if isinstance(data, dict):
                for key, value in data.items():
                    segment_text = f"Key: {key}\nValue: {json.dumps(value, indent=2)[:1000]}"
                    payload.raw_segments.append({
                        'segment_id': f"key_{key}",
                        'type': 'json_key_value',
                        'text': segment_text,
                        'metadata': {
                            'key': key,
                            'value_type': type(value).__name__,
                            'value_size': len(str(value))
                        }
                    })
            elif isinstance(data, list):
                chunk_size = min(100, len(data))  # Process in chunks
                for i in range(0, len(data), chunk_size):
                    chunk = data[i:i+chunk_size]
                    segment_text = f"Array items {i} to {min(i+chunk_size-1, len(data)-1)}:\n"
                    segment_text += json.dumps(chunk, indent=2)[:1000]
                    
                    payload.raw_segments.append({
                        'segment_id': f"array_items_{i}_to_{min(i+chunk_size-1, len(data)-1)}",
                        'type': 'json_array_chunk',
                        'text': segment_text,
                        'metadata': {
                            'start_index': i,
                            'end_index': min(i+chunk_size-1, len(data)-1),
                            'item_count': len(chunk)
                        }
                    })
            
        except Exception as e:
            logger.exception(f"JSON processing error: {e}")
            raise
    
    async def handle_text(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin1', 'cp1252']
            text_content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    text_content = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                raise ValueError("Could not decode text file with any supported encoding")
            
            payload.extracted_text = text_content
            
            # Analyze text structure
            lines = text_content.split('\n')
            paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
            
            payload.source_metadata.update({
                'encoding_used': used_encoding,
                'line_count': len(lines),
                'paragraph_count': len(paragraphs),
                'character_count': len(text_content),
                'word_count': len(text_content.split())
            })
            
            # Create structure elements for paragraphs
            char_position = 0
            for i, paragraph in enumerate(paragraphs):
                start_pos = text_content.find(paragraph, char_position)
                end_pos = start_pos + len(paragraph)
                
                payload.structure_elements.append({
                    'type': 'paragraph',
                    'paragraph_number': i + 1,
                    'length': len(paragraph),
                    'word_count': len(paragraph.split()),
                    'start_char': start_pos,
                    'end_char': end_pos
                })
                
                # Create segment for paragraph
                payload.raw_segments.append({
                    'segment_id': f"paragraph_{i + 1}",
                    'type': 'text_paragraph',
                    'text': paragraph,
                    'metadata': {
                        'paragraph_number': i + 1,
                        'word_count': len(paragraph.split()),
                        'char_start': start_pos,
                        'char_end': end_pos
                    }
                })
                
                char_position = end_pos
            
        except Exception as e:
            logger.exception(f"Text processing error: {e}")
            raise
    
    async def handle_markdown(self, file_content: bytes, payload: ParsedPayload, job: IngestJob):
        """Handle Markdown files with structure detection"""
        try:
            # Decode markdown
            encodings = ['utf-8', 'latin1', 'cp1252']
            md_content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    md_content = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if md_content is None:
                raise ValueError("Could not decode markdown file with any supported encoding")
            
            payload.extracted_text = md_content
            
            # Parse markdown structure
            lines = md_content.split('\n')
            headings = []
            current_section = None
            char_position = 0
            
            for line_num, line in enumerate(lines):
                line_stripped = line.strip()
                
                # Detect headings
                if line_stripped.startswith('#'):
                    level = len(line_stripped) - len(line_stripped.lstrip('#'))
                    heading_text = line_stripped.lstrip('# ').strip()
                    
                    heading_info = {
                        'type': 'heading',
                        'level': level,
                        'text': heading_text,
                        'line_number': line_num + 1,
                        'start_char': char_position,
                        'end_char': char_position + len(line)
                    }
                    
                    headings.append(heading_info)
                    payload.structure_elements.append(heading_info)
                    current_section = heading_text
                
                # Detect code blocks
                elif line_stripped.startswith('```'):
                    payload.structure_elements.append({
                        'type': 'code_block',
                        'line_number': line_num + 1,
                        'start_char': char_position,
                        'language': line_stripped[3:].strip() if len(line_stripped) > 3 else 'unknown'
                    })
                
                # Detect lists
                elif line_stripped.startswith(('-', '*', '+')):
                    payload.structure_elements.append({
                        'type': 'list_item',
                        'line_number': line_num + 1,
                        'start_char': char_position,
                        'text': line_stripped[1:].strip()
                    })
                
                char_position += len(line) + 1  # +1 for newline
            
            # Create segments based on sections (between headings)
            if headings:
                sections = []
                for i, heading in enumerate(headings):
                    start_char = heading['end_char']
                    end_char = headings[i + 1]['start_char'] if i + 1 < len(headings) else len(md_content)
                    
                    section_text = md_content[start_char:end_char].strip()
                    if section_text:
                        sections.append({
                            'heading': heading['text'],
                            'level': heading['level'],
                            'content': section_text,
                            'start_char': start_char,
                            'end_char': end_char
                        })
                
                # Create segments for sections
                for i, section in enumerate(sections):
                    payload.raw_segments.append({
                        'segment_id': f"section_{i + 1}",
                        'type': 'markdown_section',
                        'text': f"# {section['heading']}\n\n{section['content']}",
                        'metadata': {
                            'heading': section['heading'],
                            'heading_level': section['level'],
                            'section_number': i + 1,
                            'char_start': section['start_char'],
                            'char_end': section['end_char']
                        }
                    })
            else:
                # No headings, create paragraph-based segments
                paragraphs = [p.strip() for p in md_content.split('\n\n') if p.strip()]
                for i, paragraph in enumerate(paragraphs):
                    payload.raw_segments.append({
                        'segment_id': f"paragraph_{i + 1}",
                        'type': 'markdown_paragraph',
                        'text': paragraph,
                        'metadata': {
                            'paragraph_number': i + 1
                        }
                    })
            
            payload.source_metadata.update({
                'encoding_used': used_encoding,
                'line_count': len(lines),
                'heading_count': len(headings),
                'max_heading_level': max([h['level'] for h in headings]) if headings else 0,
                'character_count': len(md_content),
                'word_count': len(md_content.split())
            })
            
        except Exception as e:
            logger.exception(f"Markdown processing error: {e}")
            raise
    
    def _calculate_hash(self, content: bytes) -> str:
        """Calculate SHA-256 hash of content"""
        import hashlib
        return hashlib.sha256(content).hexdigest()
    
    def _parse_html_structure(self, html_content: str) -> List[Dict[str, Any]]:
        """Parse HTML structure from mammoth output"""
        # Simple HTML parsing for structure elements
        import re
        
        elements = []
        
        # Find headings
        heading_pattern = r'<h([1-6])>(.*?)</h[1-6]>'
        for match in re.finditer(heading_pattern, html_content, re.IGNORECASE):
            level = int(match.group(1))
            text = re.sub(r'<[^>]+>', '', match.group(2))  # Strip HTML tags
            elements.append({
                'type': 'heading',
                'level': level,
                'text': text.strip(),
                'html_start': match.start(),
                'html_end': match.end()
            })
        
        # Find paragraphs
        para_pattern = r'<p>(.*?)</p>'
        for match in re.finditer(para_pattern, html_content, re.IGNORECASE):
            text = re.sub(r'<[^>]+>', '', match.group(1))  # Strip HTML tags
            if text.strip():
                elements.append({
                    'type': 'paragraph',
                    'text': text.strip(),
                    'html_start': match.start(),
                    'html_end': match.end()
                })
        
        return elements
    
    def _get_json_type_description(self, data: Any) -> str:
        """Get human-readable description of JSON structure"""
        if isinstance(data, dict):
            return f"Object with {len(data)} keys"
        elif isinstance(data, list):
            return f"Array with {len(data)} items"
        elif isinstance(data, str):
            return "String value"
        elif isinstance(data, (int, float)):
            return "Numeric value"
        elif isinstance(data, bool):
            return "Boolean value"
        elif data is None:
            return "Null value"
        else:
            return f"Unknown type: {type(data).__name__}"
    
    def _analyze_json_structure(self, data: Any, path: str = "", max_depth: int = 10) -> Dict[str, Any]:
        """Recursively analyze JSON structure"""
        analysis = {
            'max_depth': 0,
            'total_keys': 0,
            'total_values': 0,
            'data_types': {},
            'elements': []
        }
        
        def analyze_recursive(obj: Any, current_path: str, depth: int):
            if depth > max_depth:
                return
            
            analysis['max_depth'] = max(analysis['max_depth'], depth)
            
            obj_type = type(obj).__name__
            analysis['data_types'][obj_type] = analysis['data_types'].get(obj_type, 0) + 1
            
            if isinstance(obj, dict):
                analysis['total_keys'] += len(obj)
                for key, value in obj.items():
                    key_path = f"{current_path}.{key}" if current_path else key
                    analysis['elements'].append({
                        'type': 'json_key',
                        'path': key_path,
                        'key': key,
                        'value_type': type(value).__name__,
                        'depth': depth
                    })
                    analyze_recursive(value, key_path, depth + 1)
            
            elif isinstance(obj, list):
                analysis['total_values'] += len(obj)
                for i, item in enumerate(obj[:10]):  # Analyze first 10 items
                    item_path = f"{current_path}[{i}]"
                    analysis['elements'].append({
                        'type': 'json_array_item',
                        'path': item_path,
                        'index': i,
                        'value_type': type(item).__name__,
                        'depth': depth
                    })
                    analyze_recursive(item, item_path, depth + 1)
            
            else:
                analysis['total_values'] += 1
        
        analyze_recursive(data, path, 0)
        return analysis
    
    def _generate_integrity_metadata(self, payload: ParsedPayload) -> Dict[str, Any]:
        """Generate integrity metadata for verification"""
        return {
            'text_length': len(payload.extracted_text),
            'segment_count': len(payload.raw_segments),
            'structure_element_count': len(payload.structure_elements),
            'concept_count': len(payload.semantic_concepts),
            'processing_complete': True,
            'extraction_method': 'production_handlers_dynamic',
            'content_hash': self._calculate_hash(payload.extracted_text.encode('utf-8')),
            'metadata_complete': bool(payload.source_metadata),
            'concept_extraction_quality': self._calculate_concept_quality_score(payload.semantic_concepts) if payload.semantic_concepts else 0.0
        }
    
    def _calculate_concept_quality_score(self, concepts: List[Dict[str, Any]]) -> float:
        """Calculate overall quality score for extracted concepts"""
        if not concepts:
            return 0.0
        
        total_confidence = sum(concept.get('confidence', 0.0) for concept in concepts)
        average_confidence = total_confidence / len(concepts)
        
        # Bonus for diverse extraction methods
        methods = set(concept.get('extraction_method', 'unknown') for concept in concepts)
        diversity_bonus = min(len(methods) * 0.1, 0.3)
        
        return min(average_confidence + diversity_bonus, 1.0)
    
    async def _extract_semantic_concepts_dynamic(self, payload: ParsedPayload):
        """
        🚀 DYNAMIC CONCEPT EXTRACTION SYSTEM 🚀
        Removes 10-concept cap and implements unlimited, quality-based extraction
        
        Features:
        - Multi-pass extraction with different algorithms
        - Dynamic quality thresholds
        - Semantic clustering and deduplication
        - Rich metadata and analytics
        - Unlimited concept extraction (quality-filtered)
        """
        try:
            extraction_start = datetime.now()
            text = payload.extracted_text
            
            logger.info(f"🔍 Starting Dynamic Concept Extraction for document {payload.document_id}")
            logger.info(f"📊 Input text: {len(text)} characters, {len(text.split())} words")
            
            # Phase 1: Multi-Algorithm Concept Extraction
            all_concepts = []
            extraction_stats = {
                'algorithms_used': [],
                'concepts_per_algorithm': {},
                'processing_time_per_algorithm': {}
            }
            
            # Algorithm 1: Enhanced Keyword Extraction
            alg_start = datetime.now()
            keyword_concepts = await self._extract_keyword_concepts(text, payload)
            alg_time = (datetime.now() - alg_start).total_seconds()
            all_concepts.extend(keyword_concepts)
            extraction_stats['algorithms_used'].append('keyword_extraction')
            extraction_stats['concepts_per_algorithm']['keyword_extraction'] = len(keyword_concepts)
            extraction_stats['processing_time_per_algorithm']['keyword_extraction'] = alg_time
            
            # Algorithm 2: Structure-Based Extraction
            alg_start = datetime.now()
            structure_concepts = await self._extract_structure_concepts(payload)
            alg_time = (datetime.now() - alg_start).total_seconds()
            all_concepts.extend(structure_concepts)
            extraction_stats['algorithms_used'].append('structure_extraction')
            extraction_stats['concepts_per_algorithm']['structure_extraction'] = len(structure_concepts)
            extraction_stats['processing_time_per_algorithm']['structure_extraction'] = alg_time
            
            # Algorithm 3: Entity Recognition
            if self.concept_config['enable_entity_extraction']:
                alg_start = datetime.now()
                entity_concepts = await self._extract_entity_concepts(text)
                alg_time = (datetime.now() - alg_start).total_seconds()
                all_concepts.extend(entity_concepts)
                extraction_stats['algorithms_used'].append('entity_extraction')
                extraction_stats['concepts_per_algorithm']['entity_extraction'] = len(entity_concepts)
                extraction_stats['processing_time_per_algorithm']['entity_extraction'] = alg_time
            
            # Algorithm 4: Frequency-Based Extraction
            if self.concept_config['enable_frequency_analysis']:
                alg_start = datetime.now()
                frequency_concepts = await self._extract_frequency_concepts(text)
                alg_time = (datetime.now() - alg_start).total_seconds()
                all_concepts.extend(frequency_concepts)
                extraction_stats['algorithms_used'].append('frequency_analysis')
                extraction_stats['concepts_per_algorithm']['frequency_analysis'] = len(frequency_concepts)
                extraction_stats['processing_time_per_algorithm']['frequency_analysis'] = alg_time
            
            logger.info(f"📈 Raw extraction complete: {len(all_concepts)} concepts from {len(extraction_stats['algorithms_used'])} algorithms")
            
            # Phase 2: Dynamic Quality Filtering
            filtered_concepts = await self._apply_dynamic_quality_filtering(all_concepts, text)
            logger.info(f"🎯 Quality filtering: {len(all_concepts)} → {len(filtered_concepts)} concepts")
            
            # Phase 3: Semantic Clustering and Deduplication
            if self.concept_config['enable_keyword_clustering']:
                clustered_concepts = await self._apply_semantic_clustering(filtered_concepts)
                logger.info(f"🔗 Semantic clustering: {len(filtered_concepts)} → {len(clustered_concepts)} concepts")
            else:
                clustered_concepts = filtered_concepts
            
            # Phase 4: Final Ranking and Selection
            final_concepts = await self._rank_and_select_concepts(clustered_concepts, text)
            
            # Phase 5: Enhance with Rich Metadata
            enhanced_concepts = await self._enhance_concept_metadata(final_concepts, text, payload)
            
            # Calculate dynamic statistics
            extraction_time = (datetime.now() - extraction_start).total_seconds()
            concept_density = len(enhanced_concepts) / max(len(text.split()), 1) * 1000  # concepts per 1000 words
            
            # Store concepts (NO CAP! 🎉)
            payload.semantic_concepts = enhanced_concepts
            
            # Enhanced logging and analytics
            logger.info(f"✅ Dynamic Concept Extraction COMPLETE!")
            logger.info(f"📊 FINAL RESULTS:")
            logger.info(f"   • Concepts extracted: {len(enhanced_concepts)} (NO CAP APPLIED)")
            logger.info(f"   • Extraction time: {extraction_time:.2f}s")
            logger.info(f"   • Concept density: {concept_density:.2f} concepts/1000 words")
            logger.info(f"   • Average confidence: {sum(c.get('confidence', 0) for c in enhanced_concepts) / max(len(enhanced_concepts), 1):.3f}")
            logger.info(f"   • Quality score: {self._calculate_concept_quality_score(enhanced_concepts):.3f}")
            
            # Log algorithm performance
            for alg, count in extraction_stats['concepts_per_algorithm'].items():
                time_taken = extraction_stats['processing_time_per_algorithm'][alg]
                logger.info(f"   • {alg}: {count} concepts in {time_taken:.3f}s")
            
            # Store extraction analytics in payload metadata
            payload.source_metadata.update({
                'concept_extraction_analytics': {
                    'total_concepts_extracted': len(enhanced_concepts),
                    'extraction_time_seconds': extraction_time,
                    'concept_density_per_1000_words': concept_density,
                    'algorithms_used': extraction_stats['algorithms_used'],
                    'concepts_per_algorithm': extraction_stats['concepts_per_algorithm'],
                    'processing_time_per_algorithm': extraction_stats['processing_time_per_algorithm'],
                    'quality_threshold_applied': self._calculate_dynamic_threshold(text),
                    'cap_removed': True,
                    'extraction_method': 'dynamic_unlimited'
                }
            })
            
        except Exception as e:
            logger.exception(f"❌ Dynamic concept extraction error: {e}")
            # Fallback to basic extraction if dynamic system fails
            payload.semantic_concepts = await self._fallback_concept_extraction(payload.extracted_text)
            logger.warning(f"⚠️ Fallback extraction used: {len(payload.semantic_concepts)} concepts")
    
    async def _extract_keyword_concepts(self, text: str, payload: ParsedPayload) -> List[Dict[str, Any]]:
        """Enhanced keyword-based concept extraction"""
        concepts = []
        text_lower = text.lower()
        
        # Expanded domain-specific concept libraries
        tech_concepts = [
            'artificial intelligence', 'machine learning', 'deep learning', 'neural network',
            'natural language processing', 'computer vision', 'data science', 'big data',
            'cloud computing', 'distributed systems', 'microservices', 'api gateway',
            'container orchestration', 'kubernetes', 'docker', 'devops', 'cicd',
            'blockchain', 'cryptocurrency', 'smart contracts', 'decentralized applications',
            'cybersecurity', 'encryption', 'authentication', 'authorization',
            'file_storage optimization', 'data warehousing', 'business intelligence',
            'predictive analytics', 'recommendation systems', 'search algorithms'
        ]
        
        business_concepts = [
            'strategic planning', 'market analysis', 'competitive advantage',
            'customer acquisition', 'retention strategy', 'revenue optimization',
            'cost reduction', 'process improvement', 'digital transformation',
            'stakeholder engagement', 'risk management', 'compliance framework',
            'performance metrics', 'key performance indicators', 'return on investment',
            'market segmentation', 'brand positioning', 'value proposition',
            'supply chain management', 'inventory optimization', 'quality assurance'
        ]
        
        academic_concepts = [
            'research methodology', 'statistical analysis', 'hypothesis testing',
            'experimental design', 'peer review', 'literature review',
            'case study analysis', 'qualitative research', 'quantitative analysis',
            'data collection', 'sampling methodology', 'variable analysis',
            'correlation analysis', 'regression modeling', 'factor analysis'
        ]
        
        all_concept_libraries = {
            'technology': tech_concepts,
            'business': business_concepts,
            'academic': academic_concepts
        }
        
        # Extract concepts from all libraries
        for domain, concept_list in all_concept_libraries.items():
            for concept in concept_list:
                if concept in text_lower:
                    # Calculate enhanced confidence based on frequency and context
                    frequency = text_lower.count(concept)
                    context_words = concept.split()
                    
                    # Boost confidence for multi-word concepts
                    confidence = min(0.7 + (len(context_words) - 1) * 0.05 + frequency * 0.02, 0.95)
                    
                    concepts.append({
                        'name': concept.title(),
                        'keywords': context_words,
                        'confidence': confidence,
                        'extraction_method': 'enhanced_keyword_matching',
                        'domain': domain,
                        'frequency': frequency,
                        'concept_length': len(context_words),
                        'source_library': domain
                    })
        
        return concepts
    
    async def _extract_structure_concepts(self, payload: ParsedPayload) -> List[Dict[str, Any]]:
        """Extract concepts from document structure elements"""
        concepts = []
        
        for element in payload.structure_elements:
            if element.get('type') == 'heading' and element.get('text'):
                heading_text = element['text'].strip()
                words = heading_text.split()
                
                # Only consider headings that could be concepts
                if self.concept_config['minimum_concept_length'] <= len(words) <= self.concept_config['maximum_concept_length']:
                    # Higher confidence for shorter, more focused headings
                    confidence = max(0.75 - (len(words) - 2) * 0.05, 0.6)
                    
                    concepts.append({
                        'name': heading_text,
                        'keywords': [word.lower() for word in words],
                        'confidence': confidence,
                        'extraction_method': 'structure_heading_analysis',
                        'source_element': 'heading',
                        'heading_level': element.get('level', 0),
                        'structure_position': element.get('start_char', 0)
                    })
            
            elif element.get('type') == 'list_item' and element.get('text'):
                item_text = element['text'].strip()
                words = item_text.split()
                
                # Extract concepts from list items (often important points)
                if self.concept_config['minimum_concept_length'] <= len(words) <= 8:  # Shorter for list items
                    concepts.append({
                        'name': item_text,
                        'keywords': [word.lower() for word in words],
                        'confidence': 0.65,
                        'extraction_method': 'structure_list_analysis',
                        'source_element': 'list_item',
                        'item_position': element.get('line_number', 0)
                    })
        
        return concepts
    
    async def _extract_entity_concepts(self, text: str) -> List[Dict[str, Any]]:
        """Simple entity extraction using pattern matching"""
        concepts = []
        
        # Pattern-based entity extraction
        patterns = {
            'organizations': r'\b[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)*\s+(?:Inc|Corp|LLC|Ltd|Company|Organization|Institute|University|College)\b',
            'technologies': r'\b(?:API|SDK|AI|ML|IoT|SaaS|PaaS|IaaS|REST|GraphQL|JSON|XML|HTTP|HTTPS|SQL|NoSQL)\b',
            'methodologies': r'\b(?:Agile|Scrum|Kanban|Waterfall|DevOps|CI/CD|TDD|BDD|MVP|POC)\b',
            'currencies': r'\$[0-9,]+(?:\.[0-9]{2})?|\b[0-9,]+\s*(?:USD|EUR|GBP|JPY|CAD|AUD)\b',
            'percentages': r'\b[0-9]+(?:\.[0-9]+)?%\b',
            'dates': r'\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+[0-9]{1,2},?\s+[0-9]{4}\b'
        }
        
        for entity_type, pattern in patterns.items():
            matches = re.finditer(pattern, text)
            for match in matches:
                entity_text = match.group().strip()
                if entity_text and len(entity_text.split()) <= 6:
                    concepts.append({
                        'name': entity_text,
                        'keywords': entity_text.lower().split(),
                        'confidence': 0.75,
                        'extraction_method': 'pattern_entity_extraction',
                        'entity_type': entity_type,
                        'pattern_matched': pattern[:50] + '...' if len(pattern) > 50 else pattern
                    })
        
        return concepts
    
    async def _extract_frequency_concepts(self, text: str) -> List[Dict[str, Any]]:
        """Extract concepts based on term frequency analysis"""
        concepts = []
        
        # Tokenize and clean text
        words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
        
        # Remove common stop words
        stop_words = {
            'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'her', 'was', 'one', 'our',
            'out', 'day', 'get', 'has', 'him', 'his', 'how', 'man', 'new', 'now', 'old', 'see', 'two', 'who',
            'its', 'said', 'each', 'make', 'most', 'over', 'such', 'very', 'what', 'with', 'have', 'from',
            'they', 'know', 'want', 'been', 'good', 'much', 'some', 'time', 'will', 'when', 'come', 'here',
            'just', 'like', 'long', 'many', 'other', 'than', 'then', 'them', 'well', 'were', 'this', 'that',
            'there', 'where', 'which', 'would', 'could', 'should', 'first', 'after', 'back', 'before',
            'through', 'during', 'about', 'between', 'into', 'within', 'without', 'under', 'above'
        }
        
        filtered_words = [word for word in words if word not in stop_words and len(word) >= 4]
        
        # Count word frequencies
        word_counts = Counter(filtered_words)
        
        # Extract high-frequency terms as potential concepts
        min_frequency = max(2, len(filtered_words) // 100)  # Dynamic threshold based on document size
        
        for word, count in word_counts.most_common(50):  # Top 50 frequent terms
            if count >= min_frequency:
                # Calculate confidence based on frequency and word characteristics
                confidence = min(0.6 + (count - min_frequency) * 0.02, 0.85)
                
                # Boost confidence for longer words (likely more specific)
                if len(word) >= 8:
                    confidence += 0.05
                
                concepts.append({
                    'name': word.title(),
                    'keywords': [word],
                    'confidence': confidence,
                    'extraction_method': 'frequency_analysis',
                    'frequency': count,
                    'relative_frequency': count / len(filtered_words),
                    'word_length': len(word)
                })
        
        # Extract frequent bigrams and trigrams
        for n in [2, 3]:
            ngrams = []
            for i in range(len(filtered_words) - n + 1):
                ngram = ' '.join(filtered_words[i:i+n])
                ngrams.append(ngram)
            
            ngram_counts = Counter(ngrams)
            min_ngram_freq = max(2, len(ngrams) // 200)
            
            for ngram, count in ngram_counts.most_common(20):
                if count >= min_ngram_freq:
                    confidence = min(0.7 + (count - min_ngram_freq) * 0.03, 0.9)
                    
                    concepts.append({
                        'name': ngram.title(),
                        'keywords': ngram.split(),
                        'confidence': confidence,
                        'extraction_method': f'frequency_analysis_{n}gram',
                        'frequency': count,
                        'ngram_size': n,
                        'relative_frequency': count / len(ngrams)
                    })
        
        return concepts
    
    async def _apply_dynamic_quality_filtering(self, concepts: List[Dict[str, Any]], text: str) -> List[Dict[str, Any]]:
        """Apply dynamic quality filtering based on document characteristics"""
        if not concepts:
            return concepts
        
        # Calculate dynamic threshold based on document characteristics
        dynamic_threshold = self._calculate_dynamic_threshold(text)
        
        # Apply confidence threshold
        filtered_concepts = [c for c in concepts if c.get('confidence', 0) >= dynamic_threshold]
        
        # Remove concepts that are too short or too long
        min_length = self.concept_config['minimum_concept_length']
        max_length = self.concept_config['maximum_concept_length']
        
        length_filtered = []
        for concept in filtered_concepts:
            name_words = len(concept.get('keywords', concept.get('name', '').split()))
            if min_length <= name_words <= max_length:
                length_filtered.append(concept)
        
        logger.info(f"🎯 Quality filtering applied:")
        logger.info(f"   • Dynamic threshold: {dynamic_threshold:.3f}")
        logger.info(f"   • Confidence filter: {len(concepts)} → {len(filtered_concepts)}")
        logger.info(f"   • Length filter: {len(filtered_concepts)} → {len(length_filtered)}")
        
        return length_filtered
    
    def _calculate_dynamic_threshold(self, text: str) -> float:
        """Calculate dynamic confidence threshold based on document characteristics"""
        base_threshold = self.concept_config['base_confidence_threshold']
        
        # Adjust based on document length
        word_count = len(text.split())
        if word_count < 500:
            # Shorter documents: lower threshold to capture more concepts
            length_adjustment = -0.1
        elif word_count > 5000:
            # Longer documents: higher threshold for quality
            length_adjustment = 0.1
        else:
            length_adjustment = 0.0
        
        # Adjust based on technical content density
        tech_indicators = ['API', 'algorithm', 'implementation', 'framework', 'protocol', 'architecture']
        tech_density = sum(1 for indicator in tech_indicators if indicator.lower() in text.lower()) / len(tech_indicators)
        tech_adjustment = tech_density * 0.05  # Up to 5% boost for technical content
        
        # Calculate final threshold
        dynamic_threshold = base_threshold + length_adjustment + tech_adjustment
        
        # Ensure threshold stays within reasonable bounds
        return max(0.5, min(0.9, dynamic_threshold))
    
    async def _apply_semantic_clustering(self, concepts: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Apply semantic clustering to group similar concepts"""
        if len(concepts) <= 1:
            return concepts
        
        # Simple keyword-based similarity clustering
        clustered_concepts = []
        processed_indices = set()
        
        for i, concept in enumerate(concepts):
            if i in processed_indices:
                continue
                
            cluster = [concept]
            concept_keywords = set(kw.lower() for kw in concept.get('keywords', []))
            
            # Find similar concepts
            for j, other_concept in enumerate(concepts[i+1:], i+1):
                if j in processed_indices:
                    continue
                    
                other_keywords = set(kw.lower() for kw in other_concept.get('keywords', []))
                
                # Calculate keyword overlap
                overlap = len(concept_keywords & other_keywords)
                similarity = overlap / max(len(concept_keywords), len(other_keywords))
                
                if similarity >= self.concept_config['similarity_threshold']:
                    cluster.append(other_concept)
                    processed_indices.add(j)
            
            # Merge cluster or keep original concept
            if len(cluster) > 1:
                # Create merged concept from cluster
                merged_concept = self._merge_concept_cluster(cluster)
                clustered_concepts.append(merged_concept)
                logger.debug(f"📎 Merged cluster of {len(cluster)} concepts: {merged_concept['name']}")
            else:
                clustered_concepts.append(concept)
            
            processed_indices.add(i)
        
        return clustered_concepts
    
    def _merge_concept_cluster(self, cluster: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Merge a cluster of similar concepts into a single concept"""
        # Use the concept with highest confidence as base
        primary_concept = max(cluster, key=lambda c: c.get('confidence', 0))
        
        # Combine keywords from all concepts in cluster
        all_keywords = set()
        extraction_methods = set()
        total_confidence = 0
        
        for concept in cluster:
            all_keywords.update(concept.get('keywords', []))
            extraction_methods.add(concept.get('extraction_method', 'unknown'))
            total_confidence += concept.get('confidence', 0)
        
        # Create merged concept
        merged_concept = primary_concept.copy()
        merged_concept.update({
            'keywords': sorted(list(all_keywords)),
            'confidence': min(total_confidence / len(cluster) + 0.05, 0.95),  # Slight boost for merged concepts
            'extraction_method': 'semantic_clustering',
            'merged_from': [c.get('name', '') for c in cluster],
            'cluster_size': len(cluster),
            'original_methods': list(extraction_methods)
        })
        
        return merged_concept
    
    async def _rank_and_select_concepts(self, concepts: List[Dict[str, Any]], text: str) -> List[Dict[str, Any]]:
        """Rank concepts and apply intelligent selection"""
        if not concepts:
            return concepts
        
        # Enhanced ranking with multiple factors
        for concept in concepts:
            score = concept.get('confidence', 0.0)
            
            # Boost score for specific extraction methods
            method = concept.get('extraction_method', '')
            if 'heading' in method:
                score += 0.1
            elif 'entity' in method:
                score += 0.08
            elif 'clustering' in method:
                score += 0.05
            
            # Boost score for multi-word concepts (more specific)
            keyword_count = len(concept.get('keywords', []))
            if keyword_count > 1:
                score += min(keyword_count * 0.02, 0.1)
            
            # Boost score for concepts with frequency data
            if 'frequency' in concept:
                frequency = concept['frequency']
                if frequency > 1:
                    score += min(frequency * 0.01, 0.05)
            
            concept['final_score'] = min(score, 1.0)
        
        # Sort by final score
        ranked_concepts = sorted(concepts, key=lambda c: c.get('final_score', 0), reverse=True)
        
        # Intelligent selection: balance quality vs quantity
        if self.concept_config['quality_over_quantity']:
            # Prefer high-quality concepts, but don't limit quantity artificially
            high_quality_threshold = self.concept_config['high_quality_threshold']
            selected_concepts = [c for c in ranked_concepts if c.get('final_score', 0) >= high_quality_threshold]
            
            # If we have very few high-quality concepts, include some medium-quality ones
            if len(selected_concepts) < 5 and len(ranked_concepts) > len(selected_concepts):
                medium_quality = [c for c in ranked_concepts if 
                                self.concept_config['base_confidence_threshold'] <= c.get('final_score', 0) < high_quality_threshold]
                selected_concepts.extend(medium_quality[:10 - len(selected_concepts)])
        else:
            # Include all concepts that meet the base threshold
            selected_concepts = ranked_concepts
        
        # Apply soft limit (can be exceeded for exceptional cases)
        max_concepts = self.concept_config['max_concepts_per_document']
        if len(selected_concepts) > max_concepts:
            # Check if we have many high-quality concepts
            exceptional_concepts = [c for c in selected_concepts[:max_concepts + 20] 
                                  if c.get('final_score', 0) >= 0.9]
            if len(exceptional_concepts) > max_concepts:
                # Allow exceeding limit for exceptional quality
                selected_concepts = exceptional_concepts
                logger.info(f"🌟 Soft limit exceeded due to exceptional concept quality: {len(exceptional_concepts)} concepts")
            else:
                selected_concepts = selected_concepts[:max_concepts]
        
        return selected_concepts
    
    async def _enhance_concept_metadata(self, concepts: List[Dict[str, Any]], text: str, payload: ParsedPayload) -> List[Dict[str, Any]]:
        """Enhance concepts with rich metadata for ψMesh integration"""
        enhanced_concepts = []
        
        for i, concept in enumerate(concepts):
            enhanced_concept = concept.copy()
            
            # Add positional information
            concept_name = concept.get('name', '').lower()
            if concept_name in text.lower():
                first_occurrence = text.lower().find(concept_name)
                enhanced_concept['first_occurrence_position'] = first_occurrence
                enhanced_concept['relative_position'] = first_occurrence / max(len(text), 1)
            
            # Add context information
            enhanced_concept.update({
                'concept_id': f"{payload.document_id}_concept_{i+1}",
                'extraction_timestamp': datetime.now().isoformat(),
                'document_context': {
                    'document_type': payload.source_metadata.get('file_type', 'unknown'),
                    'document_length': len(text),
                    'document_structure_elements': len(payload.structure_elements)
                },
                'quality_indicators': {
                    'confidence_tier': self._get_confidence_tier(concept.get('confidence', 0)),
                    'extraction_reliability': self._calculate_extraction_reliability(concept),
                    'semantic_richness': len(concept.get('keywords', [])) / max(len(concept.get('name', '').split()), 1)
                },
                'integration_metadata': {
                    'ready_for_mesh': True,
                    'concept_category': self._classify_concept_category(concept),
                    'relationship_potential': self._assess_relationship_potential(concept, concepts)
                }
            })
            
            enhanced_concepts.append(enhanced_concept)
        
        return enhanced_concepts
    
    def _get_confidence_tier(self, confidence: float) -> str:
        """Classify confidence into tiers"""
        if confidence >= 0.9:
            return 'exceptional'
        elif confidence >= 0.8:
            return 'high'
        elif confidence >= 0.7:
            return 'medium'
        elif confidence >= 0.6:
            return 'moderate'
        else:
            return 'low'
    
    def _calculate_extraction_reliability(self, concept: Dict[str, Any]) -> float:
        """Calculate reliability score based on extraction method and metadata"""
        base_reliability = concept.get('confidence', 0.0)
        
        # Boost reliability for certain extraction methods
        method = concept.get('extraction_method', '')
        if 'clustering' in method:
            base_reliability += 0.1
        elif 'heading' in method:
            base_reliability += 0.08
        elif 'entity' in method:
            base_reliability += 0.05
        
        # Boost for concepts with frequency data
        if concept.get('frequency', 0) > 1:
            base_reliability += 0.03
        
        return min(base_reliability, 1.0)
    
    def _classify_concept_category(self, concept: Dict[str, Any]) -> str:
        """Classify concept into categories for better organization"""
        method = concept.get('extraction_method', '')
        keywords = [kw.lower() for kw in concept.get('keywords', [])]
        
        # Technical concepts
        tech_indicators = ['api', 'algorithm', 'system', 'framework', 'protocol', 'architecture', 'implementation']
        if any(indicator in ' '.join(keywords) for indicator in tech_indicators):
            return 'technical'
        
        # Business concepts
        business_indicators = ['strategy', 'market', 'customer', 'revenue', 'cost', 'performance', 'management']
        if any(indicator in ' '.join(keywords) for indicator in business_indicators):
            return 'business'
        
        # Process concepts
        process_indicators = ['process', 'method', 'approach', 'procedure', 'workflow', 'operation']
        if any(indicator in ' '.join(keywords) for indicator in process_indicators):
            return 'process'
        
        # Structural concepts (from headings, etc.)
        if 'heading' in method or 'structure' in method:
            return 'structural'
        
        # Entity concepts
        if 'entity' in method:
            return 'entity'
        
        return 'general'
    
    def _assess_relationship_potential(self, concept: Dict[str, Any], all_concepts: List[Dict[str, Any]]) -> float:
        """Assess potential for relationships with other concepts"""
        concept_keywords = set(kw.lower() for kw in concept.get('keywords', []))
        
        relationship_count = 0
        for other_concept in all_concepts:
            if other_concept == concept:
                continue
                
            other_keywords = set(kw.lower() for kw in other_concept.get('keywords', []))
            overlap = len(concept_keywords & other_keywords)
            
            if overlap > 0:
                relationship_count += 1
        
        # Normalize by total concepts
        return min(relationship_count / max(len(all_concepts) - 1, 1), 1.0)
    
    async def _fallback_concept_extraction(self, text: str) -> List[Dict[str, Any]]:
        """Fallback extraction if dynamic system fails"""
        concepts = []
        
        # Simple fallback: extract common technical terms
        fallback_terms = [
            'artificial intelligence', 'machine learning', 'data analysis',
            'natural language processing', 'computer vision', 'deep learning',
            'neural network', 'algorithm', 'framework', 'architecture',
            'implementation', 'optimization', 'performance', 'system'
        ]
        
        text_lower = text.lower()
        for term in fallback_terms:
            if term in text_lower:
                concepts.append({
                    'name': term.title(),
                    'keywords': term.split(),
                    'confidence': 0.6,
                    'extraction_method': 'fallback_extraction'
                })
        
        # Limit fallback to reasonable number
        return concepts[:15]

# Global handler instance
file_handlers = ProductionFileHandlers()